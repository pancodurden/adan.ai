from __future__ import annotations

import os, io, time, re, csv, tempfile, json
from typing import Optional, List

import requests
import httpx
from fastapi import FastAPI, UploadFile, File, HTTPException, Body, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, StreamingResponse
from dotenv import load_dotenv

# ---------- .env ----------
load_dotenv()
OLLAMA_URL   = (os.getenv("OLLAMA_URL") or "http://127.0.0.1:11434").rstrip("/")
OLLAMA_MODEL = os.getenv("OLLAMA_MODEL", "deepseek-r1:7b")

ELEVEN_API_KEY  = os.getenv("sk_37c8c68595f952f1fb977e22a9af637f616ce6343c7bc30b")  # zorunlu değil; TTS çağırınca kontrol ediyoruz
ELEVEN_VOICE_ID = os.getenv("ELEVEN_VOICE_ID", "21m00Tcm4TlvDq8ikWAM")  # Rachel (varsayılan)

# ---------- Metin çıkarma kütüphaneleri ----------
# pip install pdfminer.six python-docx python-pptx openpyxl beautifulsoup4 chardet
from pdfminer.high_level import extract_text as pdf_extract_text
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation
from openpyxl import load_workbook
from bs4 import BeautifulSoup
import chardet

# ---------- Yardımcılar ----------
THINK_RE = re.compile(r"<think>[\s\S]*?</think>", re.IGNORECASE)
def strip_think(s: str) -> str:
    return THINK_RE.sub("", s or "").strip()

def _tmp_to_path(upload: UploadFile) -> str:
    """UploadFile'i geçici dosyaya yaz ve yolunu döndür."""
    suffix = os.path.splitext(upload.filename or "")[-1].lower()
    fd, path = tempfile.mkstemp(suffix=suffix)
    with os.fdopen(fd, "wb") as f:
        f.write(upload.file.read())
    upload.file.close()
    return path

def _detect_decode(raw: bytes) -> str:
    enc = chardet.detect(raw).get("encoding") or "utf-8"
    try:
        return raw.decode(enc, errors="ignore")
    except Exception:
        return raw.decode("utf-8", errors="ignore")

def extract_from_pdf(path: str) -> str:
    try:
        return pdf_extract_text(path) or ""
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"PDF parse error: {e}")

def extract_from_docx(path: str) -> str:
    try:
        doc = DocxDocument(path)
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"DOCX parse error: {e}")

def extract_from_pptx(path: str) -> str:
    try:
        pres = PptxPresentation(path)
        out: List[str] = []
        for i, slide in enumerate(pres.slides, start=1):
            texts = []
            for shp in slide.shapes:
                if hasattr(shp, "has_text_frame") and shp.has_text_frame:
                    texts.append(shp.text)
            if texts:
                out.append(f"[Slide {i}] " + "   ".join(t.strip() for t in texts if t.strip()))
        return "\n".join(out)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"PPTX parse error: {e}")

def extract_from_xlsx(path: str) -> str:
    try:
        wb = load_workbook(path, data_only=True)
        out: List[str] = []
        for ws in wb.worksheets:
            out.append(f"[Sheet] {ws.title}")
            for row in ws.iter_rows(values_only=True):
                vals = [("" if v is None else str(v)) for v in row]
                if any(v.strip() for v in vals):
                    out.append(" | ".join(vals))
        return "\n".join(out)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"XLSX parse error: {e}")

def extract_from_csv(path: str) -> str:
    try:
        raw = open(path, "rb").read()
        text = _detect_decode(raw)
        out: List[str] = []
        reader = csv.reader(io.StringIO(text))
        for row in reader:
            out.append(" | ".join(row))
        return "\n".join(out)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"CSV parse error: {e}")

def extract_from_txt_like(path: str) -> str:
    try:
        raw = open(path, "rb").read()
        return _detect_decode(raw)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Text parse error: {e}")

def extract_from_html(path: str) -> str:
    try:
        raw = open(path, "rb").read()
        html = _detect_decode(raw)
        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.extract()
        txt = soup.get_text("\n")
        return "\n".join(line.strip() for line in txt.splitlines() if line.strip())
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"HTML parse error: {e}")

def ollama_generate(prompt: str) -> str:
    """Ollama /api/generate (stream=False) tek seferlik cevap."""
    try:
        r = requests.post(
            f"{OLLAMA_URL}/api/generate",
            json={"model": OLLAMA_MODEL, "prompt": prompt, "stream": False},
            timeout=120,
        )
        r.raise_for_status()
        data = r.json()
        return strip_think((data or {}).get("response", "") or "")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Ollama error: {e}")

# ---------- FastAPI ----------
app = FastAPI(title="Adan AI Backend")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # dilersen kısıtlayabilirsin: ["http://127.0.0.1:5500"]
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ---------- Endpoints ----------
@app.get("/api/health")
def health():
    return {"ok": True, "ts": time.time(), "model": OLLAMA_MODEL, "ollama": OLLAMA_URL}

@app.post("/api/chat")
def chat(body: dict = Body(...)):
    text = (body or {}).get("text", "")
    text = text.strip() if isinstance(text, str) else ""
    if not text:
        return {"ok": True, "reply": ""}

    sys = (
        "You are Adan AI. Be concise, helpful, and reply in English. "
        "Avoid hallucinations; ask a brief question if context is missing."
    )
    prompt = f"{sys}\n\nUser: {text}\nAssistant:"
    reply = ollama_generate(prompt)
    return {"ok": True, "reply": reply, "model": OLLAMA_MODEL, "ts": time.time(), "raw": None}

@app.post("/api/extract")
async def extract(file: UploadFile = File(...)):
    if not file or not file.filename:
        raise HTTPException(status_code=400, detail="No file.")
    ext = os.path.splitext(file.filename)[-1].lower()
    path = _tmp_to_path(file)
    try:
        if ext == ".pdf":
            text = extract_from_pdf(path)
        elif ext == ".docx":
            text = extract_from_docx(path)
        elif ext == ".pptx":
            text = extract_from_pptx(path)
        elif ext in [".xlsx", ".xlsm", ".xltx"]:
            text = extract_from_xlsx(path)
        elif ext == ".csv":
            text = extract_from_csv(path)
        elif ext in [".txt", ".md", ".rtf"]:
            text = extract_from_txt_like(path)
        elif ext in [".html", ".htm"]:
            text = extract_from_html(path)
        else:
            raise HTTPException(status_code=415, detail=f"Unsupported file type: {ext}")
    finally:
        try: os.remove(path)
        except Exception: pass

    return {"ok": True, "text": text}

# ---------- ElevenLabs TTS ----------
@app.post("/api/tts")
async def eleven_tts(
    text_form: Optional[str] = Form(None),
    body: Optional[dict] = Body(None),
):
    """
    ElevenLabs TTS endpoint.
    - form-data ile:   text=...
    - JSON ile:        {"text": "..."}
    Ses MP3 (audio/mpeg) olarak stream edilir.
    """
    # gelen metni al
    text = text_form
    if not text and body and isinstance(body, dict):
        text = body.get("text")
    if not text or not isinstance(text, str) or not text.strip():
        raise HTTPException(status_code=400, detail="Missing 'text'.")

    # ortam değişkenleri kontrolü
    if not ELEVEN_API_KEY:
        raise HTTPException(status_code=500, detail="ELEVEN_API_KEY not set in environment.")

    voice_id = ELEVEN_VOICE_ID or "21m00Tcm4TlvDq8ikWAM"
    url = f"https://api.elevenlabs.io/v1/text-to-speech/{voice_id}"

    headers = {
        "Accept": "audio/mpeg",
        "Content-Type": "application/json",
        "xi-api-key": ELEVEN_API_KEY,
    }
    payload = {
        "text": text,
        # ihtiyacına göre değiştir: "eleven_multilingual_v2" (TR/EN vs. için)
        "model_id": "eleven_monolingual_v1",
        "voice_settings": {"stability": 0.6, "similarity_boost": 0.85},
    }

    async with httpx.AsyncClient(timeout=60.0) as client:
        r = await client.post(url, headers=headers, json=payload)
        if r.status_code != 200:
            # Eleven'dan gelen hata mesajını ilet
            raise HTTPException(status_code=r.status_code, detail=f"ElevenLabs error: {r.text}")

        # MP3 stream et
        return StreamingResponse(io.BytesIO(r.content), media_type="audio/mpeg")